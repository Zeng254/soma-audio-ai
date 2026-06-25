import os
import requests
import uuid
import chardet
from io import BytesIO
from typing import Literal,Callable, Any, Optional,Union
from pydantic import BaseModel, Field, field_validator,PrivateAttr,ConfigDict
from urllib.parse import urlparse
from pptx import Presentation

MAX_FILE_SIZE = 100 * 1024 * 1024

class File(BaseModel):
    """
    Generic file object, supports automatic type inference and path management
    """
    url: str = Field(..., description="FileURL(http/https)or local path")
    file_type: Literal['image', 'video', 'audio', 'document', 'default'] = Field(
        default="default",
        description="File type"
    )
    _local_path: Optional[str] = PrivateAttr(default=None)
    model_config = ConfigDict(
        json_schema_extra={
            "x-component": "file-upload",  # Frontend file upload component
        }
    )

    def set_cache_path(self, path: str):
        """Set cache path"""
        self._local_path = path

    def get_cache_path(self) -> Optional[str]:
        """Get cache path (if file actually exists)"""
        return self._local_path

    @property
    def is_remote(self) -> bool:
        """Check if network URL or local file"""
        return self.url.startswith(('http://', 'https://'))

def infer_file_category(path_or_url: str) -> tuple[str, str]:
    """
    Check file type based on path or URL suffix
    Logic:
    1. Parse URL to remove query parameters (?id=...), extract path
    2. Get last part of path as filename and suffix
    3. Lookup check, returns 'default' if no match

    Return:
        - Classification: image, video, audio, document, default
        - Suffix: .pdf

    """

    # === Step 1 & 2: Extract pure suffix ===
    # urlparse can simultaneously process local path (treated as path) and network URL
    parsed = urlparse(path_or_url)
    path = parsed.path  # Extract path part, ignore http://... and ?query=...

    # Get filename (e.g. /a/b/test.jpg -> test.jpg)
    filename = os.path.basename(path)

    # Separate suffix (test.jpg -> .jpg)
    _, ext_with_dot = os.path.splitext(filename)

    # If no suffix, fallback directly
    if not ext_with_dot:
        return 'default', ""

    # Remove dot and convert to lowercase (e.g. .JPG -> jpg)
    ext = ext_with_dot.lstrip('.').lower()

    # === Step 3: Lookup match ===
    # Define common mapping table
    TYPE_MAPPING = {
        'image': {
            'apng', 'avif', 'bmp', 'gif', 'heic', 'ico', 'jpg', 'jpeg', 'png', 'svg', 'tiff', 'webp'
        },
        'video': {
            'mp4', 'avi', 'mov', 'mkv', 'flv', 'wmv', 'webm', 'm4v', '3gp'
        },
        'audio': {
            'mp3', 'wav', 'flac', 'aac', 'ogg', 'wma', 'm4a'
        },
        'document': {
            'pdf', 'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx',
            'txt', 'md', 'csv', 'json', 'xml', 'html', 'htm'
        },
    }

    for category, extensions in TYPE_MAPPING.items():
        if ext in extensions:
            return category, ext_with_dot

    return 'default', ext_with_dot

class FileOps:
    DOWNLOAD_DIR = "/tmp"

    @staticmethod
    def _get_bytes_stream(file_obj:File) -> tuple[bytes, str]:
        """
        Get file content and suffix, size limit check, raise exception if exceeded
        """
        _, ext = infer_file_category(file_obj.url)

        if file_obj.is_remote:
            try:
                # stream=True: At this point only download headers, connection stays open, body not yet downloaded
                with requests.get(file_obj.url, stream=True, timeout=60) as resp:
                    resp.raise_for_status()

                    content_length = resp.headers.get('Content-Length')
                    if content_length and int(content_length) > MAX_FILE_SIZE:
                        raise Exception(
                            f"File size ({int(content_length)} bytes) exceeds limit 100MB, download terminated."
                        )

                    # Scenario: Header missing Content-Length or server header spoofing
                    downloaded_content = BytesIO()
                    current_size = 0

                    # Read in chunks, 8KB each
                    for chunk in resp.iter_content(chunk_size=8192):
                        if chunk:
                            current_size += len(chunk)
                            if current_size > MAX_FILE_SIZE:
                                raise Exception(f"Detected file exceeds 100MB, download interrupted.")
                            downloaded_content.write(chunk)

                    # Get complete bytes
                    return downloaded_content.getvalue(), ext

            except requests.RequestException as e:
                raise RuntimeError(f"Network request failed: {e}")

        else:
            if not os.path.exists(file_obj.url):
                raise FileNotFoundError(f"Local file does not exist: {file_obj.url}")

            '''
            file_size = os.path.getsize(file_obj.url)
            if file_size > MAX_FILE_SIZE:
                 raise Exception(f"Local file size ({file_size} bytes) exceeds limit 100MB")
            '''

            with open(file_obj.url, 'rb') as f:
                return f.read(), ext

    @staticmethod
    def save_to_local(file_obj: File, filename: str) -> str:
        """
        Save current file object content to local path, return local path
        If local path, return directly
        """
        if not file_obj.is_remote:
            if os.path.exists(file_obj.url):
                return file_obj.url

            raise FileNotFoundError(f"Local file not found: {file_obj.url}")

        try:
            os.makedirs(FileOps.DOWNLOAD_DIR, exist_ok=True)

            # Simple filename generation strategy (real scene recommends using url hash to avoid repeated downloads)
            # ext = os.path.splitext(file_obj.url.split('?')[0])[1] or ".tmp"
            # filename = f"{uuid.uuid4().hex}{ext}"
            local_path = os.path.join(FileOps.DOWNLOAD_DIR, filename)

            headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3'}
            with requests.get(file_obj.url, headers=headers, stream=True, timeout=120) as r:
                r.raise_for_status()
                with open(local_path, 'wb') as f:
                    for chunk in r.iter_content(chunk_size=8192):
                        f.write(chunk)

            return local_path
        except Exception as e:
            raise RuntimeError(f"Download failed for {file_obj.url}: {str(e)}")

    @staticmethod
    def read_bytes(file_obj:File) -> bytes:
        """
        Get file raw binary data
        Scene: Upload to OSS, save locally, pass to image processing library
        """
        content, _ = FileOps._get_bytes_stream(file_obj)
        return content

    @staticmethod
    def extract_text(file_obj: File) -> str:
        """
        ExtractionTextContent
        Scene: RAG, HTML parsing, document analysis
        """
        try:
            content, ext = FileOps._get_bytes_stream(file_obj)

            if ext in ['.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']:
                return FileOps._parse_document_bytes(file_obj, content, ext)

            # Default direct read
            charset = chardet.detect(content)
            if 'encoding' in charset:
                return content.decode(charset['encoding'])
            else:
                return content.decode('utf-8')

        except Exception as e:
            return f"[FileOps Error] Failed to read content: {str(e)}"

    @staticmethod
    def _parse_document_bytes(file_obj: File, content: bytes, ext:str) -> str:
        stream = BytesIO(content)
        text_result = ""

        try:
            if ext == '.pdf':
                import pypdf
                reader = pypdf.PdfReader(stream)
                for page in reader.pages:
                    text_result += page.extract_text() + "\n"
            elif ext in ['.docx', '.doc']:
                text_result = read_docx(stream)
            elif ext in ['.xlsx', '.xls', '.csv']:
                import pandas as pd
                if ext == '.csv':
                    df = pd.read_csv(stream)
                else:
                    df = pd.read_excel(stream)
                text_result = df.to_string()
            elif ext in ['.ppt', '.pptx']:
                text_result = read_ppt(stream)
            else:
                text_result = f"[Document format not supported for parsing: {ext}]"
        except ImportError as e:
            text_result = f"[Parsing library missing] {e}"
        except Exception as e:
            text_result = f"[ParseFail] {e}"

        return text_result

def read_docx(cont_stream) -> str:
    """
    Use docx2python to read content in order
    """
    from docx2python import docx2python
    doc_result = docx2python(cont_stream)

    # Get document structure
    all_parts = []

    # docx2python returns content as nested list
    # Iterate document body
    for section in doc_result.body:
        if isinstance(section, list):
            for item in section:
                if isinstance(item, list):
                    # May be table or multi-level content
                    for sub_item in item:
                        if isinstance(sub_item, str) and sub_item.strip():
                            all_parts.append(sub_item.strip())
                        elif isinstance(sub_item, list):
                            # Table row
                            row_text = "\n".join([str(cell).strip() for cell in sub_item if str(cell).strip()])
                            if row_text:
                                all_parts.append(row_text)
                elif isinstance(item, str) and item.strip():
                    all_parts.append(item.strip())

    # Close document
    doc_result.close()

    return "\n\n".join(all_parts)

def read_ppt(file_input: Union[str, bytes, BytesIO]) -> str:
    if not Presentation:
        return "[Error] python-pptx library not installed, cannot parse PPT file"

    # 1. Uniformly convert to file stream object (BytesIO)
    if isinstance(file_input, str):
        with open(file_input, 'rb') as f:
            ppt_stream = BytesIO(f.read())
    elif isinstance(file_input, bytes):
        ppt_stream = BytesIO(file_input)
    else:
        ppt_stream = file_input

    try:
        prs = Presentation(ppt_stream)
        full_text = []

        for i, slide in enumerate(prs.slides):
            page_content = []
            page_content.append(f"=== Page {i+1} ===")

            # shape.text_frame contains text paragraphs within shape
            for shape in slide.shapes:
                # Extract normal text boxes
                if hasattr(shape, "text") and shape.text.strip():
                    page_content.append(shape.text.strip())

                # B. Extract table content (normal shape.text cannot get text inside tables)
                if shape.has_table:
                    table_texts = []
                    for row in shape.table.rows:
                        row_cells = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()]
                        if row_cells:
                            table_texts.append(" | ".join(row_cells))
                    if table_texts:
                        page_content.append("[Table]\n" + "\n".join(table_texts))

            # Many important info hidden in notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    page_content.append(f"[Notes]: {notes.strip()}")

            full_text.append("\n".join(page_content))

        return "\n\n".join(full_text)

    except Exception as e:
        return f"[PPTParseFail] {str(e)}"


# ============================================================
# Convenience function
# ============================================================

def get_extension(path: str) -> str:
    """
    Get file extension (including dot)
    
    Args:
        path: File path
        
    Returns:
        Extension (including dot), e.g. '.wav'
    """
    import os
    return os.path.splitext(path)[1].lower()


def safe_filename(filename: str) -> str:
    """
    Generate secure filename (remove dangerous characters)
    
    Args:
        filename: Original filename
        
    Returns:
        Secure filename
    """
    import re
    # Remove or replace dangerous characters
    safe = re.sub(r'[<>:"/\\|?*\x00-\x1f]', '_', filename)
    # Remove leading and trailing spaces and dots
    safe = safe.strip('. ')
    # Limit length
    if len(safe) > 200:
        name, ext = os.path.splitext(safe)
        safe = name[:200-len(ext)] + ext
    # Ensure not empty
    if not safe:
        safe = "unnamed"
    return safe
